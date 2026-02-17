#!/usr/bin/env python3
"""
PPTX to LaTeX Beamer Converter

Converts PowerPoint presentations (.pptx) to LaTeX Beamer format using
the beamertheme.sty style file. Output defaults to beamer-slides.tex
for compatibility with beamerProject-TexStudio.txss2.

Usage:
    python pptx2tex.py [input.pptx] [output.tex]

If no arguments provided, processes all .pptx files in pptx-input/
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path
from typing import Optional

import zipfile
import shutil

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL_TYPE
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
    import io
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    print("Warning: Pillow not installed. Images will not be converted to JPG. Install with: pip install Pillow")


def patch_pptx_content_types(pptx_path: str) -> str:
    """
    Patch PPTX file to add missing content types for webp and other formats.
    Returns path to patched file (may be a temp copy).
    """
    import tempfile
    import xml.etree.ElementTree as ET

    # Create a temp copy
    temp_dir = tempfile.mkdtemp()
    temp_pptx = os.path.join(temp_dir, "patched.pptx")
    shutil.copy2(pptx_path, temp_pptx)

    # Add missing content types
    content_types_to_add = {
        '.webp': 'image/webp',
        '.svg': 'image/svg+xml',
        '.heic': 'image/heic',
        '.heif': 'image/heif',
        '.avif': 'image/avif',
        '.mp4': 'video/mp4',
        '.webm': 'video/webm',
        '.m4v': 'video/x-m4v',
        '.mov': 'video/quicktime',
    }

    try:
        with zipfile.ZipFile(temp_pptx, 'a') as zf:
            # Read existing content types
            ct_xml = zf.read('[Content_Types].xml').decode('utf-8')
            root = ET.fromstring(ct_xml)

            ns = {'ct': 'http://schemas.openxmlformats.org/package/2006/content-types'}
            ET.register_namespace('', ns['ct'])

            # Get existing extensions
            existing = set()
            for default in root.findall('.//ct:Default', ns):
                ext = default.get('Extension', '').lower()
                existing.add('.' + ext if not ext.startswith('.') else ext)

            # Add missing types
            modified = False
            for ext, content_type in content_types_to_add.items():
                if ext.lower() not in existing:
                    elem = ET.SubElement(root, '{http://schemas.openxmlformats.org/package/2006/content-types}Default')
                    elem.set('Extension', ext.lstrip('.'))
                    elem.set('ContentType', content_type)
                    modified = True

            if modified:
                # Write back
                new_ct = ET.tostring(root, encoding='unicode')
                # Remove old and add new
                with zipfile.ZipFile(pptx_path, 'r') as original:
                    with zipfile.ZipFile(temp_pptx, 'w') as new_zip:
                        for item in original.infolist():
                            if item.filename == '[Content_Types].xml':
                                new_zip.writestr(item, new_ct)
                            else:
                                new_zip.writestr(item, original.read(item.filename))

        return temp_pptx
    except Exception as e:
        print(f"Warning: Could not patch content types: {e}")
        return pptx_path


class PPTXToLatexConverter:
    """Converts PPTX presentations to LaTeX Beamer format."""

    # Standard aspect ratios (width:height ratios)
    ASPECT_RATIO_4_3 = 4 / 3      # 1.333...
    ASPECT_RATIO_16_9 = 16 / 9    # 1.777...
    ASPECT_RATIO_16_10 = 16 / 10  # 1.6

    # Beamer 16:9 dimensions (approximately)
    BEAMER_169_TEXTWIDTH_CM = 15.8   # Actual textwidth in beamer 16:9
    BEAMER_169_TEXTHEIGHT_CM = 8.0   # Usable text height (excluding header/footer)

    def __init__(self, input_path: str, output_path: str,
                 fig_dir: str = "fig", video_dir: str = "videos"):
        self.input_path = Path(input_path)
        self.output_path = Path(output_path)
        self.fig_dir = Path(fig_dir)
        self.video_dir = Path(video_dir)
        self.image_counter = 0
        self.video_counter = 0

        # Aspect ratio tracking - set during conversion
        self.source_slide_width = None   # EMU
        self.source_slide_height = None  # EMU
        self.source_aspect_ratio = None  # width/height
        self.target_aspect_ratio = self.ASPECT_RATIO_16_9  # Output is always 16:9

        # Ensure output directories exist
        self.fig_dir.mkdir(parents=True, exist_ok=True)
        self.video_dir.mkdir(parents=True, exist_ok=True)

    def detect_aspect_ratio(self, slide_width, slide_height) -> float:
        """Detect the aspect ratio of the source presentation.

        Args:
            slide_width: Width in EMU
            slide_height: Height in EMU

        Returns:
            Aspect ratio (width/height)
        """
        if not slide_width or not slide_height or slide_height == 0:
            return self.ASPECT_RATIO_4_3  # Default to 4:3

        ratio = slide_width / slide_height

        # Store for later use
        self.source_slide_width = slide_width
        self.source_slide_height = slide_height
        self.source_aspect_ratio = ratio

        return ratio

    def get_aspect_ratio_name(self, ratio: float) -> str:
        """Get a human-readable name for an aspect ratio."""
        if abs(ratio - self.ASPECT_RATIO_4_3) < 0.05:
            return "4:3"
        elif abs(ratio - self.ASPECT_RATIO_16_9) < 0.05:
            return "16:9"
        elif abs(ratio - self.ASPECT_RATIO_16_10) < 0.05:
            return "16:10"
        else:
            return f"{ratio:.2f}:1"

    def compute_canvas_dimensions(self, max_height: float = 0.75) -> tuple:
        """Compute canvas dimensions for TikZ that preserve aspect ratio relationships.

        When converting from 4:3 to 16:9, we need to handle the fact that
        the 4:3 content is relatively taller compared to width than 16:9.

        Args:
            max_height: Maximum height as fraction of textheight

        Returns:
            Tuple of (canvas_width, canvas_height) as fractions of textwidth/textheight
        """
        if self.source_aspect_ratio is None:
            # No source info - use defaults
            return (0.9, max_height)

        # Ratio of source to target aspect ratios
        # If source is 4:3 (1.333) and target is 16:9 (1.777), ratio = 0.75
        # This means source content is relatively wider for its height
        aspect_scale = self.source_aspect_ratio / self.target_aspect_ratio

        if aspect_scale < 1.0:
            # Source is narrower (taller) than target - e.g., 4:3 -> 16:9
            # Content will have letterboxing on sides
            # Scale width down to maintain proportions
            canvas_width = 0.9 * aspect_scale
            canvas_height = max_height
        else:
            # Source is wider than target (unlikely with 16:9 output)
            # Content will have letterboxing on top/bottom
            canvas_width = 0.9
            canvas_height = max_height / aspect_scale

        return (canvas_width, canvas_height)

    def transform_x_coordinate(self, rel_x: float, canvas_width: float) -> float:
        """Transform a relative X coordinate from source to target aspect ratio.

        Args:
            rel_x: X position as fraction of source slide width (0-1)
            canvas_width: Canvas width fraction

        Returns:
            X position in TikZ coordinates (fraction of textwidth)
        """
        if self.source_aspect_ratio is None:
            return rel_x * canvas_width

        aspect_scale = self.source_aspect_ratio / self.target_aspect_ratio

        if aspect_scale < 1.0:
            # Source narrower - center content horizontally
            # Add offset to center the content
            offset = (0.9 - canvas_width) / 2
            return offset + rel_x * canvas_width
        else:
            return rel_x * canvas_width

    def transform_dimensions(self, rel_w: float, rel_h: float,
                           canvas_width: float, canvas_height: float) -> tuple:
        """Transform relative dimensions from source to target aspect ratio.

        Maintains the visual proportions of shapes when converting between
        aspect ratios.

        Args:
            rel_w: Width as fraction of source slide width
            rel_h: Height as fraction of source slide height
            canvas_width: Canvas width fraction
            canvas_height: Canvas height fraction

        Returns:
            Tuple of (width_fraction, height_fraction) for TikZ
        """
        img_width = rel_w * canvas_width
        img_height = rel_h * canvas_height

        return (img_width, img_height)

    def get_textwidth_cm(self) -> float:
        """Get the effective textwidth in cm for the current aspect ratio conversion."""
        if self.source_aspect_ratio is None:
            return self.BEAMER_169_TEXTWIDTH_CM

        aspect_scale = self.source_aspect_ratio / self.target_aspect_ratio

        if aspect_scale < 1.0:
            # Content narrower - effective width is scaled
            return self.BEAMER_169_TEXTWIDTH_CM * aspect_scale
        else:
            return self.BEAMER_169_TEXTWIDTH_CM

    def escape_latex(self, text: str) -> str:
        """Escape special LaTeX characters."""
        if not text:
            return ""

        # Order matters - escape backslash first
        replacements = [
            ('\\', r'\textbackslash{}'),
            ('&', r'\&'),
            ('%', r'\%'),
            ('$', r'\$'),
            ('#', r'\#'),
            ('_', r'\_'),
            ('{', r'\{'),
            ('}', r'\}'),
            ('~', r'\textasciitilde{}'),
            ('^', r'\textasciicircum{}'),
        ]

        for old, new in replacements:
            text = text.replace(old, new)

        # Convert math symbols after escaping special characters
        text = self.convert_math_symbols(text)

        return text

    def convert_math_symbols(self, text: str) -> str:
        """Convert mathematical symbols to LaTeX math mode equivalents."""
        if not text:
            return ""

        # Unicode math symbols to LaTeX math mode
        math_symbols = [
            # Comparison operators
            ('≥', r'$\geq$'),
            ('≤', r'$\leq$'),
            ('≠', r'$\neq$'),
            ('≈', r'$\approx$'),
            ('≡', r'$\equiv$'),
            ('≢', r'$\not\equiv$'),
            ('≪', r'$\ll$'),
            ('≫', r'$\gg$'),
            ('∝', r'$\propto$'),
            ('≃', r'$\simeq$'),
            ('≅', r'$\cong$'),
            ('≲', r'$\lesssim$'),
            ('≳', r'$\gtrsim$'),
            # ASCII comparison sequences (must come after Unicode)
            ('>=', r'$\geq$'),
            ('<=', r'$\leq$'),
            ('!=', r'$\neq$'),
            ('~=', r'$\approx$'),
            ('<<', r'$\ll$'),
            ('>>', r'$\gg$'),
            # Arithmetic operators
            ('×', r'$\times$'),
            ('÷', r'$\div$'),
            ('±', r'$\pm$'),
            ('∓', r'$\mp$'),
            ('·', r'$\cdot$'),  # middle dot
            ('∙', r'$\bullet$'),
            ('⋅', r'$\cdot$'),  # dot operator
            # Greek letters (lowercase)
            ('α', r'$\alpha$'),
            ('β', r'$\beta$'),
            ('γ', r'$\gamma$'),
            ('δ', r'$\delta$'),
            ('ε', r'$\varepsilon$'),
            ('ζ', r'$\zeta$'),
            ('η', r'$\eta$'),
            ('θ', r'$\theta$'),
            ('ι', r'$\iota$'),
            ('κ', r'$\kappa$'),
            ('λ', r'$\lambda$'),
            ('μ', r'$\mu$'),
            ('ν', r'$\nu$'),
            ('ξ', r'$\xi$'),
            ('π', r'$\pi$'),
            ('ρ', r'$\rho$'),
            ('σ', r'$\sigma$'),
            ('τ', r'$\tau$'),
            ('υ', r'$\upsilon$'),
            ('φ', r'$\varphi$'),
            ('χ', r'$\chi$'),
            ('ψ', r'$\psi$'),
            ('ω', r'$\omega$'),
            # Greek letters (uppercase)
            ('Α', r'A'),  # Alpha looks like A
            ('Β', r'B'),  # Beta looks like B
            ('Γ', r'$\Gamma$'),
            ('Δ', r'$\Delta$'),
            ('Ε', r'E'),  # Epsilon looks like E
            ('Ζ', r'Z'),  # Zeta looks like Z
            ('Η', r'H'),  # Eta looks like H
            ('Θ', r'$\Theta$'),
            ('Ι', r'I'),  # Iota looks like I
            ('Κ', r'K'),  # Kappa looks like K
            ('Λ', r'$\Lambda$'),
            ('Μ', r'M'),  # Mu looks like M
            ('Ν', r'N'),  # Nu looks like N
            ('Ξ', r'$\Xi$'),
            ('Ο', r'O'),  # Omicron looks like O
            ('Π', r'$\Pi$'),
            ('Ρ', r'P'),  # Rho looks like P
            ('Σ', r'$\Sigma$'),
            ('Τ', r'T'),  # Tau looks like T
            ('Υ', r'$\Upsilon$'),
            ('Φ', r'$\Phi$'),
            ('Χ', r'X'),  # Chi looks like X
            ('Ψ', r'$\Psi$'),
            ('Ω', r'$\Omega$'),
            # Arrows
            ('→', r'$\rightarrow$'),
            ('←', r'$\leftarrow$'),
            ('↔', r'$\leftrightarrow$'),
            ('⇒', r'$\Rightarrow$'),
            ('⇐', r'$\Leftarrow$'),
            ('⇔', r'$\Leftrightarrow$'),
            ('↑', r'$\uparrow$'),
            ('↓', r'$\downarrow$'),
            ('↗', r'$\nearrow$'),
            ('↘', r'$\searrow$'),
            ('↙', r'$\swarrow$'),
            ('↖', r'$\nwarrow$'),
            ('⟶', r'$\longrightarrow$'),
            ('⟵', r'$\longleftarrow$'),
            # Set theory and logic
            ('∈', r'$\in$'),
            ('∉', r'$\notin$'),
            ('⊂', r'$\subset$'),
            ('⊃', r'$\supset$'),
            ('⊆', r'$\subseteq$'),
            ('⊇', r'$\supseteq$'),
            ('∩', r'$\cap$'),
            ('∪', r'$\cup$'),
            ('∅', r'$\emptyset$'),
            ('∧', r'$\land$'),
            ('∨', r'$\lor$'),
            ('¬', r'$\neg$'),
            ('∀', r'$\forall$'),
            ('∃', r'$\exists$'),
            ('∄', r'$\nexists$'),
            # Calculus and analysis
            ('∞', r'$\infty$'),
            ('∂', r'$\partial$'),
            ('∇', r'$\nabla$'),
            ('∫', r'$\int$'),
            ('∬', r'$\iint$'),
            ('∭', r'$\iiint$'),
            ('∮', r'$\oint$'),
            ('∑', r'$\sum$'),
            ('∏', r'$\prod$'),
            ('√', r'$\sqrt{}$'),
            ('∛', r'$\sqrt[3]{}$'),
            ('∜', r'$\sqrt[4]{}$'),
            # Miscellaneous
            ('°', r'$^\circ$'),
            ('′', r"$'$"),  # prime
            ('″', r"$''$"),  # double prime
            ('‰', r'\textperthousand{}'),
            ('…', r'\ldots{}'),
            ('ℓ', r'$\ell$'),
            ('ℏ', r'$\hbar$'),
            ('℃', r'$^\circ$C'),
            ('℉', r'$^\circ$F'),
            ('Å', r'\AA{}'),
            ('⊥', r'$\perp$'),
            ('∥', r'$\parallel$'),
            ('∠', r'$\angle$'),
            ('△', r'$\triangle$'),
            ('□', r'$\square$'),
            ('◇', r'$\diamond$'),
            ('★', r'$\star$'),
            ('☆', r'$\star$'),
            ('✓', r'$\checkmark$'),
            ('✗', r'$\times$'),
            ('†', r'$\dagger$'),
            ('‡', r'$\ddagger$'),
            ('§', r'\S{}'),
            ('¶', r'\P{}'),
            ('©', r'\copyright{}'),
            ('®', r'\textregistered{}'),
            ('™', r'\texttrademark{}'),
            # Superscripts and subscripts (common ones)
            ('²', r'$^2$'),
            ('³', r'$^3$'),
            ('¹', r'$^1$'),
            ('⁰', r'$^0$'),
            ('⁴', r'$^4$'),
            ('⁵', r'$^5$'),
            ('⁶', r'$^6$'),
            ('⁷', r'$^7$'),
            ('⁸', r'$^8$'),
            ('⁹', r'$^9$'),
            ('⁺', r'$^+$'),
            ('⁻', r'$^-$'),
            ('₀', r'$_0$'),
            ('₁', r'$_1$'),
            ('₂', r'$_2$'),
            ('₃', r'$_3$'),
            ('₄', r'$_4$'),
            ('₅', r'$_5$'),
            ('₆', r'$_6$'),
            ('₇', r'$_7$'),
            ('₈', r'$_8$'),
            ('₉', r'$_9$'),
            ('₊', r'$_+$'),
            ('₋', r'$_-$'),
            # Fractions
            ('½', r'$\frac{1}{2}$'),
            ('⅓', r'$\frac{1}{3}$'),
            ('⅔', r'$\frac{2}{3}$'),
            ('¼', r'$\frac{1}{4}$'),
            ('¾', r'$\frac{3}{4}$'),
            ('⅕', r'$\frac{1}{5}$'),
            ('⅖', r'$\frac{2}{5}$'),
            ('⅗', r'$\frac{3}{5}$'),
            ('⅘', r'$\frac{4}{5}$'),
            ('⅙', r'$\frac{1}{6}$'),
            ('⅚', r'$\frac{5}{6}$'),
            ('⅛', r'$\frac{1}{8}$'),
            ('⅜', r'$\frac{3}{8}$'),
            ('⅝', r'$\frac{5}{8}$'),
            ('⅞', r'$\frac{7}{8}$'),
        ]

        for old, new in math_symbols:
            text = text.replace(old, new)

        return text

    def clean_text(self, text: str) -> str:
        """Clean and normalize text from PPTX."""
        if not text:
            return ""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        return text

    def rgb_to_latex_color(self, rgb_color) -> Optional[str]:
        """Convert python-pptx RGB color to LaTeX color definition."""
        if rgb_color is None:
            return None
        try:
            # rgb_color is an RGBColor object with red, green, blue attributes (0-255)
            r, g, b = rgb_color[0], rgb_color[1], rgb_color[2]
            # Return as RGB values for \textcolor[RGB]{r,g,b}{text}
            return f"{r},{g},{b}"
        except (TypeError, IndexError, AttributeError):
            return None

    def get_font_size_command(self, font_size_pt: float, base_size_pt: float = 11.0) -> Optional[str]:
        """Convert font size to LaTeX size command relative to base size."""
        if font_size_pt is None or font_size_pt <= 0:
            return None

        ratio = font_size_pt / base_size_pt

        # Map ratio to LaTeX size commands
        if ratio < 0.6:
            return "\\tiny"
        elif ratio < 0.75:
            return "\\scriptsize"
        elif ratio < 0.85:
            return "\\footnotesize"
        elif ratio < 0.95:
            return "\\small"
        elif ratio < 1.1:
            return None  # Normal size, no command needed
        elif ratio < 1.3:
            return "\\large"
        elif ratio < 1.6:
            return "\\Large"
        elif ratio < 2.0:
            return "\\LARGE"
        elif ratio < 2.5:
            return "\\huge"
        else:
            return "\\Huge"

    def format_run_to_latex(self, run, base_font_size: float = 11.0, para_font=None) -> str:
        """Convert a single text run to LaTeX with formatting.

        Args:
            run: The text run from python-pptx
            base_font_size: Base font size for relative sizing
            para_font: Paragraph-level font as fallback for inherited properties
        """
        text = run.text
        if not text:
            return ""

        # Escape LaTeX special characters
        escaped_text = self.escape_latex(text)

        # Track formatting wrappers
        prefix = ""
        suffix = ""

        # Check for font properties
        font = run.font

        # Bold - check run font, then paragraph font
        is_bold = font.bold
        if is_bold is None and para_font:
            is_bold = para_font.bold
        if is_bold:
            prefix += "\\textbf{"
            suffix = "}" + suffix

        # Italic - check run font, then paragraph font
        is_italic = font.italic
        if is_italic is None and para_font:
            is_italic = para_font.italic
        if is_italic:
            prefix += "\\textit{"
            suffix = "}" + suffix

        # Font size - check run font, then paragraph font
        font_size = None
        try:
            if font.size is not None:
                font_size = font.size.pt
            elif para_font and para_font.size is not None:
                font_size = para_font.size.pt

            if font_size is not None:
                size_cmd = self.get_font_size_command(font_size, base_font_size)
                if size_cmd:
                    prefix = "{" + size_cmd + " " + prefix
                    suffix = suffix + "}"
        except (AttributeError, TypeError):
            pass

        # Font color
        try:
            rgb = None
            # Check run-level color
            if font.color and font.color.type is not None:
                if font.color.rgb:
                    rgb = self.rgb_to_latex_color(font.color.rgb)
                elif font.color.theme_color is not None:
                    # Theme colors - map common ones
                    # This is a simplified mapping; actual theme colors depend on the template
                    pass

            if rgb and rgb != "0,0,0":
                prefix = f"\\textcolor[RGB]{{{rgb}}}" + "{" + prefix
                suffix = suffix + "}"
        except (AttributeError, TypeError):
            pass

        return prefix + escaped_text + suffix

    def needs_space_between(self, prev_text: str, curr_text: str) -> bool:
        """Determine if a space is needed between two text segments."""
        if not prev_text or not curr_text:
            return False

        prev_char = prev_text[-1]
        curr_char = curr_text[0]

        # Already has spacing
        if prev_char in ' \n\t' or curr_char in ' \n\t':
            return False

        # Don't add space after opening brackets or before closing ones
        if prev_char in '([{' or curr_char in ')]}':
            return False

        # Don't add space around colons that are part of times or ratios
        if prev_char == ':' and curr_char.isdigit():
            return False

        # Add space between: word ending and capital letter starting (new word)
        if prev_char.isalpha() and curr_char.isupper():
            return True

        # Add space between: lowercase/digit and uppercase
        if (prev_char.islower() or prev_char.isdigit()) and curr_char.isupper():
            return True

        # Add space between: letter and digit (except subscripts like H2O)
        if prev_char.isalpha() and curr_char.isdigit():
            # Check if it looks like a unit (m³, m², etc.) - no space needed
            if prev_text.endswith(('m', 'cm', 'km', 's', 'h', 'kg', 'g', 'l', 'L')):
                return False
            return True

        # Add space between: digit and letter (like "77.9m" -> "77.9 m")
        if prev_char.isdigit() and curr_char.isalpha():
            return True

        # Add space between: punctuation (except specific ones) and alphanumeric
        if prev_char in '.!?;' and curr_char.isalnum():
            return True

        return False

    def process_paragraph_with_formatting(self, para, base_font_size: float = 11.0) -> dict:
        """Process a paragraph preserving run-level formatting."""
        formatted_parts = []
        raw_text_parts = []

        # Get paragraph-level font as fallback for inherited properties
        para_font = None
        try:
            para_font = para.font
        except AttributeError:
            pass

        # Check if paragraph has runs
        runs = list(para.runs)

        if runs:
            for run in runs:
                if run.text:
                    formatted_parts.append(self.format_run_to_latex(run, base_font_size, para_font))
                    raw_text_parts.append(run.text)
        else:
            # Fallback: if no runs, use paragraph text directly
            if para.text:
                escaped = self.escape_latex(para.text)
                formatted_parts.append(escaped)
                raw_text_parts.append(para.text)

        # Join with space where needed
        formatted_text = ""
        for i, part in enumerate(formatted_parts):
            if i > 0:
                prev_raw = raw_text_parts[i-1]
                curr_raw = raw_text_parts[i]
                if self.needs_space_between(prev_raw, curr_raw):
                    formatted_text += " "
            formatted_text += part

        level = para.level if para.level is not None else 0

        return {
            'text': formatted_text,
            'raw_text': ''.join(raw_text_parts),
            'level': level,
            'is_bullet': para.level is not None and para.level >= 0
        }

    def extract_image(self, shape, slide_num: int, slide_width=None, slide_height=None) -> Optional[dict]:
        """Extract image from shape, convert to high-resolution JPG (220 dpi), and save to fig directory.

        Returns dict with 'filename', 'width_ratio', and bounding box info for overlay detection.
        """
        try:
            if hasattr(shape, 'image'):
                image = shape.image
                self.image_counter += 1
                filename = f"slide{slide_num:02d}_img{self.image_counter:03d}.jpg"
                filepath = self.fig_dir / filename

                if HAS_PIL:
                    # Convert to JPG at 220 dpi
                    img = Image.open(io.BytesIO(image.blob))
                    # Convert to RGB if necessary (e.g., PNG with transparency)
                    if img.mode in ('RGBA', 'LA', 'P'):
                        # Create white background for transparent images
                        background = Image.new('RGB', img.size, (255, 255, 255))
                        if img.mode == 'P':
                            img = img.convert('RGBA')
                        background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                        img = background
                    elif img.mode != 'RGB':
                        img = img.convert('RGB')
                    # Save as JPG with 220 dpi
                    img.save(filepath, 'JPEG', quality=95, dpi=(220, 220))
                else:
                    # Fallback: save original format if PIL not available
                    ext = image.ext
                    filename = f"slide{slide_num:02d}_img{self.image_counter:03d}.{ext}"
                    filepath = self.fig_dir / filename
                    with open(filepath, 'wb') as f:
                        f.write(image.blob)

                # Calculate width ratio relative to slide
                width_ratio = 0.8  # Default
                if slide_width and hasattr(shape, 'width') and shape.width:
                    width_ratio = min(0.95, shape.width / slide_width)
                    # Round to reasonable precision
                    width_ratio = round(width_ratio, 2)

                # Store bounding box in EMU for overlay detection
                return {
                    'filename': filename,
                    'width_ratio': width_ratio,
                    'left': shape.left if hasattr(shape, 'left') else 0,
                    'top': shape.top if hasattr(shape, 'top') else 0,
                    'width': shape.width if hasattr(shape, 'width') else 0,
                    'height': shape.height if hasattr(shape, 'height') else 0,
                }
        except Exception as e:
            print(f"Warning: Could not extract image: {e}")
        return None

    def extract_video(self, shape, slide_num: int) -> Optional[str]:
        """Extract video from shape and save to videos directory."""
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                # Access the media part
                if hasattr(shape, '_element'):
                    # Try to get video from the shape's media
                    for rel in shape.part.rels.values():
                        if "video" in rel.reltype.lower():
                            self.video_counter += 1
                            # Determine extension from content type
                            content_type = getattr(rel.target_part, 'content_type', 'video/mp4')
                            ext = content_type.split('/')[-1] if '/' in content_type else 'mp4'
                            filename = f"slide{slide_num:02d}_vid{self.video_counter:03d}.{ext}"
                            filepath = self.video_dir / filename

                            with open(filepath, 'wb') as f:
                                f.write(rel.target_part.blob)

                            return filename
        except Exception as e:
            print(f"Warning: Could not extract video: {e}")
        return None

    def process_text_frame(self, text_frame, base_font_size: float = 11.0) -> list:
        """Process text frame and return list of paragraphs with formatting info."""
        paragraphs = []
        for para in text_frame.paragraphs:
            # Use the new formatting-aware processor
            para_info = self.process_paragraph_with_formatting(para, base_font_size)
            # Only include if there's actual content
            if para_info['raw_text'].strip():
                paragraphs.append(para_info)
        return paragraphs

    def emu_to_cm(self, emu) -> float:
        """Convert EMU (English Metric Units) to centimeters."""
        if emu is None:
            return 0.0
        # 1 inch = 914400 EMU, 1 inch = 2.54 cm
        return emu / 914400 * 2.54

    def emu_to_textwidth(self, emu, slide_width) -> float:
        """Convert EMU to fraction of textwidth."""
        if emu is None or slide_width is None or slide_width == 0:
            return 0.8
        return min(0.95, emu / slide_width)

    def shapes_overlap(self, shape, image_info: dict, tolerance: float = 0.1) -> bool:
        """Check if a shape overlaps with an image bounding box.

        Args:
            shape: The shape to check
            image_info: Dict with 'left', 'top', 'width', 'height' in EMU
            tolerance: Overlap tolerance as fraction of image size (0.1 = 10%)

        Returns True if the shape's center is within the image bounds (with tolerance).
        """
        try:
            # Get shape bounds
            shape_left = shape.left if hasattr(shape, 'left') else 0
            shape_top = shape.top if hasattr(shape, 'top') else 0
            shape_width = shape.width if hasattr(shape, 'width') else 0
            shape_height = shape.height if hasattr(shape, 'height') else 0

            # Calculate shape center
            shape_center_x = shape_left + shape_width / 2
            shape_center_y = shape_top + shape_height / 2

            # Get image bounds with tolerance
            img_left = image_info['left']
            img_top = image_info['top']
            img_width = image_info['width']
            img_height = image_info['height']

            # Expand image bounds by tolerance
            tol_x = img_width * tolerance
            tol_y = img_height * tolerance

            # Check if shape center is within expanded image bounds
            in_x = (img_left - tol_x) <= shape_center_x <= (img_left + img_width + tol_x)
            in_y = (img_top - tol_y) <= shape_center_y <= (img_top + img_height + tol_y)

            return in_x and in_y
        except (AttributeError, TypeError, KeyError):
            return False

    def get_shape_fill_color(self, shape) -> Optional[tuple]:
        """Extract fill color from shape as (r, g, b, opacity) tuple.

        Returns None if no fill or transparent.
        """
        try:
            fill = shape.fill
            if fill is None or fill.type is None:
                return None

            # Check for solid fill
            if fill.type == MSO_FILL_TYPE.SOLID:
                fore_color = fill.fore_color
                if fore_color and fore_color.rgb:
                    r, g, b = fore_color.rgb[0], fore_color.rgb[1], fore_color.rgb[2]
                    # Get transparency (0 = opaque, 1 = transparent)
                    # In python-pptx, transparency is stored differently
                    opacity = 1.0
                    try:
                        # Some shapes have transparency attribute
                        if hasattr(fill, '_fill') and hasattr(fill._fill, 'attrib'):
                            # Check for alpha in the fill
                            pass
                    except Exception:
                        pass
                    return (r, g, b, opacity)
        except Exception:
            pass
        return None

    def get_shape_font_info(self, shape) -> dict:
        """Extract font information from a shape's text.

        Returns dict with 'size_pt', 'bold', 'italic', 'color_rgb'.
        """
        result = {'size_pt': None, 'bold': False, 'italic': False, 'color_rgb': None}
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        if font.size:
                            result['size_pt'] = font.size.pt
                        if font.bold:
                            result['bold'] = True
                        if font.italic:
                            result['italic'] = True
                        if font.color and font.color.rgb:
                            result['color_rgb'] = (
                                font.color.rgb[0],
                                font.color.rgb[1],
                                font.color.rgb[2]
                            )
                        # Use first run's properties
                        if result['size_pt']:
                            break
                    if result['size_pt']:
                        break
        except Exception:
            pass
        return result

    def text_shape_to_overlay_info(self, shape, image_info: dict) -> Optional[dict]:
        """Extract overlay information from a text shape overlapping an image.

        Returns dict with 'text', 'rel_x', 'rel_y', 'has_bold', 'has_italic',
        'font_size_pt', 'font_color_rgb', 'fill_color' for rendering.
        """
        try:
            # Get shape text with formatting preserved
            shape_text = ""
            has_bold = False
            has_italic = False
            font_size_pt = None
            font_color_rgb = None

            if hasattr(shape, 'text_frame') and shape.text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        # Check paragraph-level formatting
                        para_font = para.font if hasattr(para, 'font') else None

                        # Process runs for formatting
                        for run in para.runs:
                            if run.text.strip():
                                font = run.font
                                # Check bold - run level or paragraph level
                                if font.bold:
                                    has_bold = True
                                elif para_font and para_font.bold:
                                    has_bold = True
                                # Check italic
                                if font.italic:
                                    has_italic = True
                                elif para_font and para_font.italic:
                                    has_italic = True
                                # Get font size
                                if font.size and font_size_pt is None:
                                    font_size_pt = font.size.pt
                                elif para_font and para_font.size and font_size_pt is None:
                                    font_size_pt = para_font.size.pt
                                # Get font color (handle theme colors gracefully)
                                try:
                                    if font.color and font.color.type is not None and font_color_rgb is None:
                                        if font.color.rgb:
                                            font_color_rgb = (
                                                font.color.rgb[0],
                                                font.color.rgb[1],
                                                font.color.rgb[2]
                                            )
                                except (AttributeError, TypeError):
                                    pass

                        shape_text += self.escape_latex(para.text.strip()) + " "
                shape_text = shape_text.strip()

            if not shape_text:
                return None

            # Get shape position relative to image
            shape_left = shape.left if hasattr(shape, 'left') else 0
            shape_top = shape.top if hasattr(shape, 'top') else 0

            img_left = image_info['left']
            img_top = image_info['top']
            img_width = image_info['width']
            img_height = image_info['height']

            # Calculate position as fraction of image dimensions
            # x: 0 = left edge, 1 = right edge
            # y: 0 = top edge, 1 = bottom edge
            rel_x = (shape_left - img_left) / img_width if img_width else 0.5
            rel_y = (shape_top - img_top) / img_height if img_height else 0.5

            # Clamp to valid range
            rel_x = max(0, min(1, rel_x))
            rel_y = max(0, min(1, rel_y))

            # Get fill color
            fill_color = self.get_shape_fill_color(shape)

            return {
                'text': shape_text,
                'rel_x': rel_x,
                'rel_y': rel_y,
                'has_bold': has_bold,
                'has_italic': has_italic,
                'font_size_pt': font_size_pt,
                'font_color_rgb': font_color_rgb,
                'fill_color': fill_color,
            }

        except Exception as e:
            print(f"Warning: Could not extract overlay info: {e}")
            return None

    def render_image_with_overlays(self, img_info: dict, width_ratio: float, height_ratio: float) -> list:
        """Render an image with its overlays in a TikZ picture.

        Args:
            img_info: Image info dict with 'filename' and optional 'overlays' list
            width_ratio: Width as fraction of textwidth
            height_ratio: Height as fraction of textheight

        Returns list of LaTeX lines.
        """
        lines = []
        filename = img_info['filename']
        overlays = img_info.get('overlays', [])

        if not overlays:
            # No overlays - simple includegraphics
            lines.append(f"\\begin{{center}}")
            lines.append(f"  \\includegraphics[width={width_ratio}\\textwidth,height={height_ratio}\\textheight,keepaspectratio]{{{filename}}}")
            lines.append(f"\\end{{center}}")
        else:
            # Has overlays - use TikZ with image as node and overlays positioned relative to it
            lines.append(f"\\begin{{center}}")
            lines.append(f"\\begin{{tikzpicture}}")
            lines.append(f"  \\node[anchor=center, inner sep=0] (img) {{\\includegraphics[width={width_ratio}\\textwidth,height={height_ratio}\\textheight,keepaspectratio]{{{filename}}}}};")

            for overlay in overlays:
                text = overlay['text']
                rel_x = overlay['rel_x']
                rel_y = overlay['rel_y']
                has_bold = overlay.get('has_bold', False)
                has_italic = overlay.get('has_italic', False)
                fill_color = overlay.get('fill_color')

                # Build node style
                style_parts = []

                # Background fill
                if fill_color:
                    r, g, b, opacity = fill_color
                    style_parts.append(f"fill={{rgb,255:red,{r};green,{g};blue,{b}}}")
                    if opacity < 1.0:
                        style_parts.append(f"fill opacity={opacity:.2f}")
                else:
                    # Default semi-transparent white background for readability
                    style_parts.append("fill=white")
                    style_parts.append("fill opacity=0.8")

                style_parts.append("draw=none")
                style_parts.append("rounded corners=2pt")
                style_parts.append("inner sep=2pt")
                style_parts.append("anchor=north west")

                style_str = ", ".join(style_parts)

                # Text formatting
                font_cmds = []
                if has_bold:
                    font_cmds.append("\\bfseries")
                if has_italic:
                    font_cmds.append("\\itshape")
                font_prefix = " ".join(font_cmds) + " " if font_cmds else ""

                # Position relative to image node using TikZ calc library
                # rel_x: 0 = left edge, 1 = right edge
                # rel_y: 0 = top edge, 1 = bottom edge
                # Use two-step interpolation:
                # 1. Interpolate horizontally from north west to north east
                # 2. Then interpolate vertically from that point down to south
                # The syntax $(A)!factor!(B)$ gives a point factor of the way from A to B
                lines.append(f"  \\node[{style_str}] at ($($(img.north west)!{rel_x:.3f}!(img.north east)$)!{rel_y:.3f}!($(img.south west)!{rel_x:.3f}!(img.south east)$)$) {{{font_prefix}{text}}};")

            lines.append(f"\\end{{tikzpicture}}")
            lines.append(f"\\end{{center}}")

        return lines

    def text_shape_to_tikz_overlay(self, shape, image_info: dict, slide_width, slide_height) -> Optional[str]:
        """Convert a text shape overlapping an image to a TikZ overlay node.

        Positions the text relative to the image and preserves styling.
        NOTE: This method is deprecated - use text_shape_to_overlay_info instead.
        """
        try:
            # Get shape text with formatting preserved
            shape_text = ""
            has_bold = False
            has_italic = False
            font_size_pt = None
            font_color_rgb = None

            if hasattr(shape, 'text_frame') and shape.text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        # Check paragraph-level formatting
                        para_font = para.font if hasattr(para, 'font') else None

                        # Process runs for formatting
                        for run in para.runs:
                            if run.text.strip():
                                font = run.font
                                # Check bold - run level or paragraph level
                                if font.bold:
                                    has_bold = True
                                elif para_font and para_font.bold:
                                    has_bold = True
                                # Check italic
                                if font.italic:
                                    has_italic = True
                                elif para_font and para_font.italic:
                                    has_italic = True
                                # Get font size
                                if font.size and font_size_pt is None:
                                    font_size_pt = font.size.pt
                                elif para_font and para_font.size and font_size_pt is None:
                                    font_size_pt = para_font.size.pt
                                # Get font color (handle theme colors gracefully)
                                try:
                                    if font.color and font.color.type is not None and font_color_rgb is None:
                                        if font.color.rgb:
                                            font_color_rgb = (
                                                font.color.rgb[0],
                                                font.color.rgb[1],
                                                font.color.rgb[2]
                                            )
                                except (AttributeError, TypeError):
                                    pass

                        shape_text += self.escape_latex(para.text.strip()) + " "
                shape_text = shape_text.strip()

            if not shape_text:
                return None

            # Get shape position relative to image
            shape_left = shape.left if hasattr(shape, 'left') else 0
            shape_top = shape.top if hasattr(shape, 'top') else 0
            shape_width = shape.width if hasattr(shape, 'width') else 0

            img_left = image_info['left']
            img_top = image_info['top']
            img_width = image_info['width']
            img_height = image_info['height']

            # Calculate position as fraction of image dimensions
            # x: 0 = left edge, 1 = right edge
            # y: 0 = top edge, 1 = bottom edge
            rel_x = (shape_left - img_left) / img_width if img_width else 0.5
            rel_y = (shape_top - img_top) / img_height if img_height else 0.5

            # Clamp to valid range
            rel_x = max(0, min(1, rel_x))
            rel_y = max(0, min(1, rel_y))

            # Convert to TikZ positioning (relative to image center)
            # Image width in textwidth units
            img_width_ratio = image_info.get('width_ratio', 0.8)

            # Calculate offset from image center in textwidth units
            x_offset = (rel_x - 0.5) * img_width_ratio
            # Y offset: positive = up in TikZ, but we want to go down from top
            # Estimate image height based on aspect ratio (assume ~0.6 of width for typical images)
            img_height_ratio = img_width_ratio * 0.6  # Approximation
            y_offset = (0.5 - rel_y) * img_height_ratio * 10  # Scale for cm

            # Get shape fill color
            fill_color = self.get_shape_fill_color(shape)

            # Build TikZ style
            style_parts = []

            # Background fill
            if fill_color:
                r, g, b, opacity = fill_color
                # Define color inline
                style_parts.append(f"fill={{rgb,255:red,{r};green,{g};blue,{b}}}")
                if opacity < 1.0:
                    style_parts.append(f"fill opacity={opacity:.2f}")
            else:
                # Default semi-transparent white background for readability
                style_parts.append("fill=white")
                style_parts.append("fill opacity=0.8")

            # Shape outline
            style_parts.append("draw=none")
            style_parts.append("rounded corners=2pt")

            # Inner padding
            style_parts.append("inner sep=2pt")

            # Text formatting commands
            font_cmds = []
            if has_bold:
                font_cmds.append("\\bfseries")
            if has_italic:
                font_cmds.append("\\itshape")
            if font_size_pt:
                if font_size_pt < 8:
                    font_cmds.append("\\tiny")
                elif font_size_pt < 10:
                    font_cmds.append("\\scriptsize")
                elif font_size_pt < 11:
                    font_cmds.append("\\footnotesize")
                elif font_size_pt < 12:
                    font_cmds.append("\\small")
                elif font_size_pt > 18:
                    font_cmds.append("\\Large")
                elif font_size_pt > 14:
                    font_cmds.append("\\large")

            font_prefix = " ".join(font_cmds) + " " if font_cmds else ""

            # Text color
            text_color_cmd = ""
            if font_color_rgb:
                r, g, b = font_color_rgb
                text_color_cmd = f"\\color[RGB]{{{r},{g},{b}}}"

            style_str = ", ".join(style_parts)

            # Build TikZ code - position relative to the previous image
            # Using scope with shift based on calculated position
            tikz_code = f"""\\begin{{tikzpicture}}[overlay, remember picture]
  \\node[{style_str}, anchor=north west] at ([xshift={x_offset:.3f}\\textwidth, yshift={y_offset:.2f}cm]current page.center) {{{text_color_cmd}{font_prefix}{shape_text}}};
\\end{{tikzpicture}}"""

            return tikz_code

        except Exception as e:
            print(f"Warning: Could not convert text shape to TikZ overlay: {e}")
            return None

    def extract_all_images_from_shape(self, shape, slide_num: int, slide_width=None, slide_height=None, depth=0) -> list:
        """Recursively extract images from any shape type.

        Handles: Pictures, grouped shapes, placeholder images, fill images.
        Returns list of dicts with 'filename', 'width_ratio', 'left', 'top'.
        """
        images = []
        max_depth = 5  # Prevent infinite recursion

        if depth > max_depth:
            return images

        try:
            # Handle grouped shapes - recurse into children
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child_shape in shape.shapes:
                    images.extend(self.extract_all_images_from_shape(
                        child_shape, slide_num, slide_width, slide_height, depth + 1
                    ))
                return images

            # Handle regular pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_info = self.extract_image(shape, slide_num, slide_width, slide_height)
                if img_info:
                    # Add position info
                    img_info['left'] = self.emu_to_textwidth(shape.left, slide_width) if slide_width else 0
                    img_info['top'] = self.emu_to_cm(shape.top)
                    images.append(img_info)
                return images

            # Handle placeholder shapes that might contain images
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                # Check if placeholder has an image
                if hasattr(shape, 'image'):
                    img_info = self.extract_image(shape, slide_num, slide_width, slide_height)
                    if img_info:
                        img_info['left'] = self.emu_to_textwidth(shape.left, slide_width) if slide_width else 0
                        img_info['top'] = self.emu_to_cm(shape.top)
                        images.append(img_info)
                return images

            # Handle shapes with fill images (background images on shapes)
            if hasattr(shape, 'fill'):
                try:
                    fill = shape.fill
                    if fill.type is not None and hasattr(fill, 'picture') and fill.picture:
                        # Extract the fill image
                        self.image_counter += 1
                        filename = f"slide{slide_num:02d}_img{self.image_counter:03d}.jpg"
                        filepath = self.fig_dir / filename

                        if HAS_PIL and hasattr(fill, '_fill'):
                            # Try to get the image blob from the fill
                            blob = fill._fill.blip.blob
                            if blob:
                                img = Image.open(io.BytesIO(blob))
                                if img.mode in ('RGBA', 'LA', 'P'):
                                    background = Image.new('RGB', img.size, (255, 255, 255))
                                    if img.mode == 'P':
                                        img = img.convert('RGBA')
                                    background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                                    img = background
                                elif img.mode != 'RGB':
                                    img = img.convert('RGB')
                                img.save(filepath, 'JPEG', quality=95, dpi=(220, 220))

                                width_ratio = self.emu_to_textwidth(shape.width, slide_width)
                                images.append({
                                    'filename': filename,
                                    'width_ratio': round(width_ratio, 2),
                                    'left': self.emu_to_textwidth(shape.left, slide_width) if slide_width else 0,
                                    'top': self.emu_to_cm(shape.top)
                                })
                except Exception:
                    pass  # Fill image extraction is optional

        except Exception as e:
            print(f"Warning: Could not extract image from shape: {e}")

        return images

    def shape_to_tikz(self, shape, slide_width, slide_height) -> Optional[str]:
        """Convert an AutoShape to TikZ code.

        Returns TikZ code string or None if not an AutoShape with content.
        """
        try:
            # Only handle AutoShapes
            if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                return None

            # Get shape text
            shape_text = ""
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        shape_text += self.escape_latex(para.text.strip()) + " "
                shape_text = shape_text.strip()

            if not shape_text:
                return None  # Skip shapes without text

            # Get position and size in cm
            left_cm = self.emu_to_cm(shape.left)
            top_cm = self.emu_to_cm(shape.top)
            width_cm = self.emu_to_cm(shape.width)
            height_cm = self.emu_to_cm(shape.height)

            # Convert slide coordinates to TikZ coordinates
            # Slide coordinates: origin top-left, y increases downward
            # TikZ coordinates: we'll use origin top-left too for simplicity
            slide_width_cm = self.emu_to_cm(slide_width) if slide_width else 25.4  # ~10 inches default
            slide_height_cm = self.emu_to_cm(slide_height) if slide_height else 19.05  # 7.5 inches for 4:3 aspect

            # Get canvas dimensions based on aspect ratio
            canvas_width, canvas_height = self.compute_canvas_dimensions(0.75)

            # Normalize to textwidth/textheight-based coordinates
            rel_x = left_cm / slide_width_cm
            rel_y = 1.0 - (top_cm / slide_height_cm)  # Flip y for TikZ

            # Transform for aspect ratio
            x_pos = self.transform_x_coordinate(rel_x, canvas_width)
            y_pos = rel_y * canvas_height

            # Adjust text width for aspect ratio
            effective_textwidth = self.get_textwidth_cm()
            adj_width_cm = (width_cm / slide_width_cm) * effective_textwidth

            # Determine shape style based on AutoShape type
            shape_style = "draw, rounded corners, fill=yellow!20"

            # Check if it's a callout type
            try:
                auto_shape_type = shape.auto_shape_type
                if auto_shape_type is not None:
                    type_name = str(auto_shape_type).lower()
                    if 'callout' in type_name:
                        shape_style = "draw, fill=yellow!30, rounded corners, drop shadow"
                    elif 'arrow' in type_name:
                        shape_style = "draw, ->, thick"
                    elif 'rectangle' in type_name or 'rect' in type_name:
                        shape_style = "draw, fill=blue!10"
                    elif 'oval' in type_name or 'ellipse' in type_name:
                        shape_style = "draw, ellipse, fill=green!10"
            except Exception:
                pass

            # Generate TikZ code using textheight-based positioning (consistent with render_unified_tikz)
            tikz_code = f"""\\begin{{tikzpicture}}[overlay, remember picture]
  \\node[{shape_style}, text width={adj_width_cm:.1f}cm, align=center, anchor=north west] at ([xshift={x_pos:.3f}\\textwidth, yshift={y_pos:.3f}\\textheight]current page.south west) {{{shape_text}}};
\\end{{tikzpicture}}"""

            return tikz_code

        except Exception as e:
            print(f"Warning: Could not convert shape to TikZ: {e}")
            return None

    def get_shape_line_color(self, shape) -> Optional[tuple]:
        """Extract line/border color from shape as (r, g, b) tuple.

        Returns None if no line or transparent.
        """
        try:
            line = shape.line
            if line is None:
                return None

            # Check if line has a fill (visible border)
            if line.fill and line.fill.type is not None:
                if line.fill.type == MSO_FILL_TYPE.SOLID:
                    fore_color = line.fill.fore_color
                    if fore_color and fore_color.rgb:
                        return (fore_color.rgb[0], fore_color.rgb[1], fore_color.rgb[2])

            # Alternative: check line color directly
            if hasattr(line, 'color') and line.color and line.color.rgb:
                return (line.color.rgb[0], line.color.rgb[1], line.color.rgb[2])

        except Exception:
            pass
        return None

    def classify_shape_type(self, shape) -> dict:
        """Classify a shape into a category for TikZ rendering.

        Returns dict with 'category' (rectangle, rounded_rect, oval, callout, arrow, other),
        'has_border', 'border_color', 'tikz_shape_style'.
        """
        result = {
            'category': 'rectangle',
            'has_border': True,
            'border_color': None,
            'tikz_shape_style': ''
        }

        try:
            if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                result['category'] = 'text'
                result['has_border'] = False
                return result

            auto_type = shape.auto_shape_type
            if auto_type is None:
                return result

            type_name = str(auto_type).upper()

            # Classify based on shape type name
            if 'OVAL' in type_name or 'ELLIPSE' in type_name or 'CIRCLE' in type_name:
                result['category'] = 'oval'
                result['tikz_shape_style'] = 'ellipse'
            elif 'ROUNDED' in type_name:
                result['category'] = 'rounded_rect'
                result['tikz_shape_style'] = 'rounded corners=4pt'
            elif 'CALLOUT' in type_name:
                result['category'] = 'callout'
                result['tikz_shape_style'] = 'rounded corners, drop shadow'
            elif 'ARROW' in type_name:
                result['category'] = 'arrow'
                result['tikz_shape_style'] = '->, thick'
            elif 'DIAMOND' in type_name:
                result['category'] = 'diamond'
                result['tikz_shape_style'] = 'diamond'
            elif 'TRIANGLE' in type_name:
                result['category'] = 'triangle'
                result['tikz_shape_style'] = 'regular polygon, regular polygon sides=3'
            elif 'PENTAGON' in type_name:
                result['category'] = 'pentagon'
                result['tikz_shape_style'] = 'regular polygon, regular polygon sides=5'
            elif 'HEXAGON' in type_name:
                result['category'] = 'hexagon'
                result['tikz_shape_style'] = 'regular polygon, regular polygon sides=6'
            elif 'STAR' in type_name:
                result['category'] = 'star'
                result['tikz_shape_style'] = 'star'
            elif 'CLOUD' in type_name:
                result['category'] = 'cloud'
                result['tikz_shape_style'] = 'cloud, cloud puffs=10'
            else:
                # Default to rectangle
                result['category'] = 'rectangle'

            # Get border color
            result['border_color'] = self.get_shape_line_color(shape)

            # Check if shape has visible border
            try:
                if hasattr(shape, 'line') and shape.line:
                    # Line width of 0 or no fill means no border
                    if shape.line.width and shape.line.width > 0:
                        result['has_border'] = True
                    else:
                        result['has_border'] = False
            except Exception:
                pass

        except Exception:
            pass

        return result

    def extract_positioned_element(self, shape, slide_width, slide_height) -> Optional[dict]:
        """Extract a positioned element (text shape or AutoShape) with its position and styling.

        Returns dict with 'type', 'text', 'rel_x', 'rel_y', 'rel_width', 'rel_height',
        'has_bold', 'has_italic', 'font_size_pt', 'font_color_rgb', 'fill_color',
        'auto_shape_type', 'shape_category', 'border_color', 'tikz_shape_style'.
        """
        try:
            # Get shape position relative to slide
            shape_left = shape.left if hasattr(shape, 'left') else 0
            shape_top = shape.top if hasattr(shape, 'top') else 0
            shape_width = shape.width if hasattr(shape, 'width') else 0
            shape_height = shape.height if hasattr(shape, 'height') else 0

            # Calculate relative positions (0-1 range)
            rel_x = shape_left / slide_width if slide_width else 0.5
            rel_y = shape_top / slide_height if slide_height else 0.5
            rel_width = shape_width / slide_width if slide_width else 0.1
            rel_height = shape_height / slide_height if slide_height else 0.1

            # Classify the shape type
            shape_info = self.classify_shape_type(shape)

            # Extract text and formatting
            shape_text = ""
            has_bold = False
            has_italic = False
            font_size_pt = None
            font_color_rgb = None

            if hasattr(shape, 'text_frame') and shape.text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        para_font = para.font if hasattr(para, 'font') else None
                        for run in para.runs:
                            if run.text.strip():
                                font = run.font
                                if font.bold:
                                    has_bold = True
                                elif para_font and para_font.bold:
                                    has_bold = True
                                if font.italic:
                                    has_italic = True
                                elif para_font and para_font.italic:
                                    has_italic = True
                                if font.size and font_size_pt is None:
                                    font_size_pt = font.size.pt
                                elif para_font and para_font.size and font_size_pt is None:
                                    font_size_pt = para_font.size.pt
                                try:
                                    if font.color and font.color.type is not None and font_color_rgb is None:
                                        if font.color.rgb:
                                            font_color_rgb = (font.color.rgb[0], font.color.rgb[1], font.color.rgb[2])
                                except (AttributeError, TypeError):
                                    pass
                        shape_text += self.escape_latex(para.text.strip()) + " "
                shape_text = shape_text.strip()

            if not shape_text:
                return None

            # Get fill color with better extraction
            fill_color = self.get_shape_fill_color(shape)

            # Get AutoShape type if applicable
            auto_shape_type = None
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    auto_shape_type = str(shape.auto_shape_type) if shape.auto_shape_type else None
            except Exception:
                pass

            return {
                'type': 'text' if auto_shape_type is None else 'autoshape',
                'text': shape_text,
                'rel_x': rel_x,
                'rel_y': rel_y,
                'rel_width': rel_width,
                'rel_height': rel_height,
                'has_bold': has_bold,
                'has_italic': has_italic,
                'font_size_pt': font_size_pt,
                'font_color_rgb': font_color_rgb,
                'fill_color': fill_color,
                'shape_category': shape_info['category'],
                'border_color': shape_info['border_color'],
                'has_border': shape_info['has_border'],
                'tikz_shape_style': shape_info['tikz_shape_style'],
                'auto_shape_type': auto_shape_type,
            }

        except Exception as e:
            print(f"Warning: Could not extract positioned element: {e}")
            return None

    def get_all_shape_content(self, slide, slide_num: int, slide_width, slide_height) -> dict:
        """Extract all content from a slide: title, text items, images, shapes.

        Uses a comprehensive approach to detect:
        - Images with their bounding boxes
        - Positioned elements (text shapes and AutoShapes that should maintain position)
        - Body placeholder content (bullet points)

        Returns dict with 'title', 'subtitle', 'content_items', 'images', 'videos',
        'positioned_elements', 'needs_unified_tikz'.
        """
        result = {
            'title': '',
            'subtitle': '',
            'content_items': [],
            'images': [],
            'videos': [],
            'positioned_elements': []  # New: stores all positioned shapes
        }

        # Track if we've found title/subtitle
        title_found = False
        subtitle_found = False

        # Collect shapes for two-pass processing
        all_shapes = list(slide.shapes)
        text_shapes_to_process = []  # Will be processed in second pass

        # === PASS 1: Extract all images first (needed for overlap detection) ===
        for shape in all_shapes:
            try:
                shape_type = shape.shape_type

                # Handle grouped shapes - extract images recursively
                if shape_type == MSO_SHAPE_TYPE.GROUP:
                    result['images'].extend(
                        self.extract_all_images_from_shape(shape, slide_num, slide_width, slide_height)
                    )
                    # Queue grouped text shapes for second pass
                    for child_shape in shape.shapes:
                        if hasattr(child_shape, 'text_frame') and child_shape.text:
                            text_shapes_to_process.append(child_shape)
                    continue

                # Check for image in ANY shape that has an 'image' attribute
                if hasattr(shape, 'image') and shape.image is not None:
                    img_info = self.extract_image(shape, slide_num, slide_width, slide_height)
                    if img_info:
                        result['images'].append(img_info)
                    if shape_type == MSO_SHAPE_TYPE.PICTURE:
                        continue

                # Handle pictures specifically
                if shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_info = self.extract_image(shape, slide_num, slide_width, slide_height)
                    if img_info:
                        result['images'].append(img_info)
                    continue

                # Handle videos/media
                if shape_type == MSO_SHAPE_TYPE.MEDIA:
                    vid_file = self.extract_video(shape, slide_num)
                    if vid_file:
                        result['videos'].append(vid_file)
                    continue

                # Queue other shapes for second pass
                text_shapes_to_process.append(shape)

            except Exception as e:
                print(f"Warning: Could not process shape in pass 1: {e}")
                continue

        # === PASS 2: Process text shapes ===
        for shape in text_shapes_to_process:
            try:
                shape_type = shape.shape_type

                # Handle text frames
                if hasattr(shape, 'text_frame') and shape.text and shape.text.strip():
                    # Check if this is a placeholder (title, subtitle, body)
                    is_placeholder_handled = False
                    is_body_placeholder = False
                    try:
                        if shape.is_placeholder:
                            ph_type = shape.placeholder_format.type
                            # Title placeholder
                            if ph_type in [1, 3] and not title_found:  # TITLE or CENTER_TITLE
                                result['title'] = self.clean_text(shape.text)
                                title_found = True
                                is_placeholder_handled = True
                            # Subtitle placeholder (type 2)
                            elif ph_type == 2 and not subtitle_found:
                                result['subtitle'] = self.clean_text(shape.text)
                                subtitle_found = True
                                is_placeholder_handled = True
                            # Body/Content placeholder - process as content (bullet points)
                            elif ph_type in [6, 7]:  # BODY, OBJECT
                                paragraphs = self.process_text_frame(shape.text_frame)
                                result['content_items'].extend(paragraphs)
                                is_placeholder_handled = True
                                is_body_placeholder = True
                    except (ValueError, AttributeError):
                        pass

                    # Skip source/citation text (only if it's a short citation, not body content)
                    if not is_placeholder_handled:
                        shape_text = shape.text.strip()
                        shape_text_lower = shape_text.lower()
                        # Only treat as source if it's short (single line citation) and starts with source pattern
                        is_short_text = len(shape_text) < 200 and '\n' not in shape_text
                        starts_with_source = any(shape_text_lower.startswith(pattern) for pattern in [
                            'quelle:', 'source:', 'bildquelle:', 'datenquelle:',
                            'videoquelle:', 'grafikquelle:'
                        ])
                        if is_short_text and starts_with_source:
                            # Extract source text for later use
                            result.setdefault('sources', []).append(shape_text)
                            is_placeholder_handled = True

                    # Non-placeholder text shapes: determine if they should be positioned elements
                    # or content items based on their characteristics
                    if not is_placeholder_handled and not is_body_placeholder:
                        shape_text = shape.text.strip()
                        # Heuristic: if text has multiple paragraphs/lines or is very long,
                        # it's likely body content, not a positioned label
                        num_paras = len([p for p in shape.text_frame.paragraphs if p.text.strip()])
                        is_substantial_content = num_paras > 2 or len(shape_text) > 300

                        if is_substantial_content:
                            # Treat as body content (bullet points)
                            paragraphs = self.process_text_frame(shape.text_frame)
                            result['content_items'].extend(paragraphs)
                        else:
                            # Treat as positioned element
                            elem = self.extract_positioned_element(shape, slide_width, slide_height)
                            if elem:
                                result['positioned_elements'].append(elem)
                        continue

                # Handle AutoShapes (callouts, etc.) - always positioned elements
                elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    elem = self.extract_positioned_element(shape, slide_width, slide_height)
                    if elem:
                        result['positioned_elements'].append(elem)

            except Exception as e:
                print(f"Warning: Could not process shape in pass 2: {e}")
                continue

        # Determine if we need unified TikZ rendering
        # Criteria: multiple images OR (images + positioned elements)
        num_images = len(result['images'])
        num_positioned = len(result['positioned_elements'])
        result['needs_unified_tikz'] = (num_images >= 2) or (num_images >= 1 and num_positioned >= 1)

        return result

    def is_title_slide(self, slide, slide_index: int = 0) -> bool:
        """Detect if slide is likely a title slide."""
        layout_name = slide.slide_layout.name.lower() if slide.slide_layout.name else ""

        # Check layout name for title indicators
        if 'title' in layout_name and 'content' not in layout_name:
            return True

        # First slide with certain characteristics is likely a title slide
        if slide_index == 0:
            # Count text shapes and check for typical title slide patterns
            text_shapes = [s for s in slide.shapes if hasattr(s, 'text') and s.text.strip()]
            # Title slides typically have few text elements (title, subtitle, author)
            if len(text_shapes) <= 4:
                return True

        return False

    def is_thank_you_slide(self, slide) -> bool:
        """Detect if slide is likely a thank-you slide."""
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_lower = shape.text.lower()
                if any(phrase in text_lower for phrase in ['thank you', 'thanks', 'danke', 'vielen dank']):
                    return True
        return False

    def extract_title_info(self, slide) -> dict:
        """Extract title, subtitle, and author from title slide.

        Note: In PPTX files, the SUBTITLE placeholder typically contains the
        author name, while the BODY placeholder contains the chapter/subtitle text.
        We map these to LaTeX fields accordingly:
        - PPTX SUBTITLE placeholder -> LaTeX \\author{}
        - PPTX BODY placeholder -> LaTeX \\subtitle{}
        """
        info = {'title': '', 'subtitle': '', 'author': '', 'date': ''}
        other_texts = []  # Collect non-placeholder text

        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text = self.clean_text(shape.text)
                if not text:
                    continue

                # Try to identify by placeholder type
                try:
                    if shape.is_placeholder:
                        ph_type = shape.placeholder_format.type
                        # Title placeholder types
                        if ph_type in [1, 3]:  # TITLE or CENTER_TITLE
                            info['title'] = text
                        elif ph_type == 2:  # PPTX SUBTITLE contains author name
                            info['author'] = text
                        elif ph_type == 6:  # BODY contains chapter/subtitle
                            if not info['subtitle']:
                                info['subtitle'] = text
                        elif ph_type == 13:  # DATE
                            info['date'] = text
                        continue
                except (ValueError, AttributeError):
                    pass

                # Collect other text shapes for potential subtitle detection
                other_texts.append(text)

        # If no subtitle found in placeholders, use first non-title/author text
        if not info['subtitle'] and other_texts:
            for text in other_texts:
                if text != info['title'] and text != info['author']:
                    info['subtitle'] = text
                    break

        if not info['title'] and other_texts:
            # First non-empty text as title if no title found yet
            info['title'] = other_texts[0]

        return info

    def extract_sources_from_content(self, content_items: list) -> tuple:
        """Extract source references from content items.

        Returns tuple of (filtered_content_items, sources_list).
        Sources are identified by patterns like "Bildquelle:", "Quelle:", URLs, etc.
        """
        filtered_items = []
        sources = []

        source_patterns = [
            r'(?:Bild|Daten|Video|Grafik)?[Qq]uelle:?\s*(.+)',
            r'Source:?\s*(.+)',
        ]

        for item in content_items:
            raw_text = item.get('raw_text', item.get('text', ''))
            is_source = False

            # Check for source patterns
            for pattern in source_patterns:
                match = re.match(pattern, raw_text.strip(), re.IGNORECASE)
                if match:
                    source_text = match.group(1).strip()
                    if source_text:
                        sources.append(source_text)
                    is_source = True
                    break

            # Check for standalone URLs that look like sources
            if not is_source:
                url_match = re.match(r'^https?://\S+$', raw_text.strip())
                if url_match:
                    sources.append(raw_text.strip())
                    is_source = True

            if not is_source:
                filtered_items.append(item)

        return filtered_items, sources

    def images_overlap(self, img1: dict, img2: dict, threshold: float = 0.3) -> bool:
        """Check if two images overlap significantly.

        Args:
            img1, img2: Image info dicts with 'left', 'top', 'width', 'height' in EMU
            threshold: Minimum overlap ratio to consider as overlapping (0.3 = 30%)

        Returns True if the images overlap by at least the threshold amount.
        """
        try:
            # Get bounding boxes
            l1, t1 = img1.get('left', 0), img1.get('top', 0)
            w1, h1 = img1.get('width', 0), img1.get('height', 0)
            l2, t2 = img2.get('left', 0), img2.get('top', 0)
            w2, h2 = img2.get('width', 0), img2.get('height', 0)

            # Calculate overlap region
            overlap_left = max(l1, l2)
            overlap_top = max(t1, t2)
            overlap_right = min(l1 + w1, l2 + w2)
            overlap_bottom = min(t1 + h1, t2 + h2)

            # Check if there's actual overlap
            if overlap_right <= overlap_left or overlap_bottom <= overlap_top:
                return False

            # Calculate overlap area
            overlap_area = (overlap_right - overlap_left) * (overlap_bottom - overlap_top)

            # Calculate areas of both images
            area1 = w1 * h1
            area2 = w2 * h2
            smaller_area = min(area1, area2) if area1 > 0 and area2 > 0 else 1

            # Check if overlap is significant relative to smaller image
            overlap_ratio = overlap_area / smaller_area
            return overlap_ratio >= threshold

        except (TypeError, KeyError):
            return False

    def group_overlapping_images(self, images: list) -> list:
        """Group overlapping images together for animation.

        Returns list of image groups. Each group is a list of images that overlap.
        Non-overlapping images are in their own single-element groups.
        Images within a group are ordered by their z-order (first = back, last = front).
        """
        if not images:
            return []

        # Track which images have been assigned to groups
        assigned = [False] * len(images)
        groups = []

        for i, img in enumerate(images):
            if assigned[i]:
                continue

            # Start a new group with this image
            group = [img]
            assigned[i] = True

            # Find all images that overlap with any image in the group
            changed = True
            while changed:
                changed = False
                for j, other_img in enumerate(images):
                    if assigned[j]:
                        continue
                    # Check if other_img overlaps with any image in the group
                    for group_img in group:
                        if self.images_overlap(group_img, other_img):
                            group.append(other_img)
                            assigned[j] = True
                            changed = True
                            break

            groups.append(group)

        return groups

    def estimate_content_density(self, content_items: list, images: list, videos: list,
                                   positioned_elements: list) -> dict:
        """Estimate content density to determine if adjustments are needed.

        Returns dict with 'item_count', 'image_count', 'is_text_heavy', 'is_image_heavy',
        'has_overlapping_images', 'image_groups', 'recommended_font_size', 'recommended_image_height'.
        """
        item_count = len(content_items)
        image_count = len(images)
        positioned_count = len(positioned_elements)

        # Group overlapping images
        image_groups = self.group_overlapping_images(images)
        has_overlapping_images = any(len(group) > 1 for group in image_groups)

        # Count distinct image positions (groups, not individual overlapping images)
        distinct_image_positions = len(image_groups)

        # Text-heavy: many items with few images
        is_text_heavy = item_count > 12 and image_count <= 1

        # Image-heavy: multiple non-overlapping images that need to share vertical space
        # Overlapping images don't count as image-heavy since they stack
        is_image_heavy = distinct_image_positions >= 2 and not has_overlapping_images

        # Determine recommended font size for text-heavy slides
        # Use scriptsize as the smallest to maintain readability
        recommended_font_size = None
        recommended_itemsep = None
        if is_text_heavy:
            if item_count > 18:
                recommended_font_size = "\\scriptsize"
                recommended_itemsep = "0pt"
            elif item_count > 14:
                recommended_font_size = "\\footnotesize"
                recommended_itemsep = "1pt"
            elif item_count > 10:
                recommended_font_size = "\\small"
                recommended_itemsep = "2pt"

        # Determine recommended image height for multi-image slides (non-overlapping only)
        recommended_image_height = None
        if is_image_heavy:
            available_height = 0.70  # fraction of textheight
            # Subtract space for text items if present
            if item_count > 0:
                available_height -= min(0.2, item_count * 0.02)
            # Divide among distinct image positions
            per_image_height = available_height / distinct_image_positions
            # Ensure minimum reasonable size
            per_image_height = max(0.15, min(0.5, per_image_height))
            recommended_image_height = per_image_height

        return {
            'item_count': item_count,
            'image_count': image_count,
            'is_text_heavy': is_text_heavy,
            'is_image_heavy': is_image_heavy,
            'has_overlapping_images': has_overlapping_images,
            'image_groups': image_groups,
            'recommended_font_size': recommended_font_size,
            'recommended_itemsep': recommended_itemsep,
            'recommended_image_height': recommended_image_height,
        }

    def render_unified_tikz(self, images: list, positioned_elements: list, max_height: float = 0.75) -> list:
        """Render all images and positioned elements in a unified TikZ picture.

        All elements are positioned relative to each other using their original
        slide coordinates, maintaining correct relative positions when scaled.
        Overlapping images are rendered with animation (click-to-reveal).

        Args:
            images: List of image info dicts with bounding boxes
            positioned_elements: List of positioned element dicts
            max_height: Maximum height as fraction of textheight

        Returns list of LaTeX lines.
        """
        lines = []

        if not images and not positioned_elements:
            return lines

        # Group overlapping images for animation
        image_groups = self.group_overlapping_images(images)
        has_overlapping = any(len(g) > 1 for g in image_groups)

        # Calculate the bounding box of all elements to determine scale
        all_elements = []

        for img in images:
            all_elements.append({
                'type': 'image',
                'left': img.get('left', 0),
                'top': img.get('top', 0),
                'width': img.get('width', 0),
                'height': img.get('height', 0),
                'data': img
            })

        if not all_elements and not positioned_elements:
            return lines

        # Find bounding box of all images (in EMU)
        if all_elements:
            min_left = min(e['left'] for e in all_elements)
            min_top = min(e['top'] for e in all_elements)
            max_right = max(e['left'] + e['width'] for e in all_elements)
            max_bottom = max(e['top'] + e['height'] for e in all_elements)
        else:
            min_left = min_top = 0
            max_right = max_bottom = 1

        total_width = max_right - min_left if max_right > min_left else 1
        total_height = max_bottom - min_top if max_bottom > min_top else 1

        lines.append("\\begin{center}")
        lines.append("\\begin{tikzpicture}")

        # Determine the TikZ canvas size based on content and aspect ratio
        # This properly handles 4:3 to 16:9 conversion by adjusting canvas proportions
        canvas_width, canvas_height = self.compute_canvas_dimensions(max_height)

        # Render images - handle overlapping groups with animation
        img_node_idx = 0
        for group in image_groups:
            if len(group) == 1:
                # Single image - render normally
                img = group[0]
                elem = next(e for e in all_elements if e['data'] is img)
                filename = img['filename']

                rel_x = (elem['left'] - min_left) / total_width if total_width else 0.5
                rel_y = (elem['top'] - min_top) / total_height if total_height else 0.5
                rel_w = elem['width'] / total_width if total_width else 1
                rel_h = elem['height'] / total_height if total_height else 1

                # Transform coordinates with aspect ratio compensation
                x_pos = self.transform_x_coordinate(rel_x, canvas_width)
                y_pos = (1 - rel_y - rel_h) * canvas_height
                img_width, img_height = self.transform_dimensions(rel_w, rel_h, canvas_width, canvas_height)

                lines.append(f"  \\node[anchor=south west, inner sep=0] (img{img_node_idx}) at ({x_pos:.3f}\\textwidth, {y_pos:.3f}\\textheight) {{\\includegraphics[width={img_width:.3f}\\textwidth,height={img_height:.3f}\\textheight,keepaspectratio]{{{filename}}}}};")

                # Render overlays for this image
                for overlay in img.get('overlays', []):
                    self._render_overlay_node(lines, overlay, f"img{img_node_idx}")

                img_node_idx += 1
            else:
                # Multiple overlapping images - use animation
                # Find the combined bounding box for this group
                group_elems = [next(e for e in all_elements if e['data'] is img) for img in group]
                group_left = min(e['left'] for e in group_elems)
                group_top = min(e['top'] for e in group_elems)
                group_right = max(e['left'] + e['width'] for e in group_elems)
                group_bottom = max(e['top'] + e['height'] for e in group_elems)

                # Position for the group (use the first/background image position)
                rel_x = (group_left - min_left) / total_width if total_width else 0.5
                rel_y = (group_top - min_top) / total_height if total_height else 0.5
                group_w = (group_right - group_left) / total_width if total_width else 1
                group_h = (group_bottom - group_top) / total_height if total_height else 1

                # Transform coordinates with aspect ratio compensation
                x_pos = self.transform_x_coordinate(rel_x, canvas_width)
                y_pos = (1 - rel_y - group_h) * canvas_height
                img_width, img_height = self.transform_dimensions(group_w, group_h, canvas_width, canvas_height)

                base_node = f"img{img_node_idx}"

                for anim_idx, img in enumerate(group):
                    filename = img['filename']
                    width_ratio = img.get('width_ratio', 0.8)

                    if anim_idx == 0:
                        # First image (background) - always visible
                        lines.append(f"  \\node[anchor=south west, inner sep=0] ({base_node}) at ({x_pos:.3f}\\textwidth, {y_pos:.3f}\\textheight) {{\\includegraphics[width={img_width:.3f}\\textwidth,height={img_height:.3f}\\textheight,keepaspectratio]{{{filename}}}}};")
                    else:
                        # Subsequent images - use onslide for animation
                        overlay_num = anim_idx + 1
                        lines.append(f"  \\onslide<{overlay_num}->{{\\node[anchor=south west, inner sep=0] at ({base_node}.south west) {{\\includegraphics[width={img_width:.3f}\\textwidth,height={img_height:.3f}\\textheight,keepaspectratio]{{{filename}}}}};}}")

                    # Render overlays for this image
                    for overlay in img.get('overlays', []):
                        if anim_idx == 0:
                            self._render_overlay_node(lines, overlay, base_node)
                        else:
                            # Overlays on animated images also need onslide
                            overlay_num = anim_idx + 1
                            # For now, skip overlays on non-base animated images
                            pass

                img_node_idx += 1

        # Render positioned elements (text shapes, AutoShapes)
        for elem in positioned_elements:
            self._render_positioned_element(lines, elem, canvas_width, canvas_height)

        lines.append("\\end{tikzpicture}")
        lines.append("\\end{center}")

        return lines

    def _render_overlay_node(self, lines: list, overlay: dict, parent_node: str):
        """Render an overlay as a TikZ node positioned relative to parent."""
        text = overlay['text']
        rel_x = overlay['rel_x']
        rel_y = overlay['rel_y']
        has_bold = overlay.get('has_bold', False)
        has_italic = overlay.get('has_italic', False)
        fill_color = overlay.get('fill_color')

        # Build node style
        style_parts = []
        if fill_color:
            r, g, b, opacity = fill_color
            style_parts.append(f"fill={{rgb,255:red,{r};green,{g};blue,{b}}}")
            if opacity < 1.0:
                style_parts.append(f"fill opacity={opacity:.2f}")
        else:
            style_parts.append("fill=white")
            style_parts.append("fill opacity=0.8")

        style_parts.extend(["draw=none", "rounded corners=2pt", "inner sep=2pt", "anchor=north west"])
        style_str = ", ".join(style_parts)

        # Text formatting
        font_cmds = []
        if has_bold:
            font_cmds.append("\\bfseries")
        if has_italic:
            font_cmds.append("\\itshape")
        font_prefix = " ".join(font_cmds) + " " if font_cmds else ""

        # Position relative to parent node
        lines.append(f"  \\node[{style_str}] at ($({parent_node}.north west)!{rel_x:.3f}!({parent_node}.north east)!{rel_y:.3f}!($({parent_node}.south west)!{rel_x:.3f}!({parent_node}.south east)$)$) {{{font_prefix}{text}}};")

    def _render_positioned_element(self, lines: list, elem: dict, canvas_width: float, canvas_height: float):
        """Render a positioned element (text shape or AutoShape) as a TikZ node."""
        text = elem['text']
        rel_x = elem['rel_x']
        rel_y = elem['rel_y']
        rel_width = elem.get('rel_width', 0.1)
        has_bold = elem.get('has_bold', False)
        has_italic = elem.get('has_italic', False)
        font_size_pt = elem.get('font_size_pt')
        font_color_rgb = elem.get('font_color_rgb')
        fill_color = elem.get('fill_color')
        shape_category = elem.get('shape_category', 'rectangle')
        border_color = elem.get('border_color')
        has_border = elem.get('has_border', True)
        tikz_shape_style = elem.get('tikz_shape_style', '')

        # Build node style based on element type and shape category
        style_parts = []

        # Handle shape-specific styles
        if shape_category == 'oval':
            style_parts.append("ellipse")
        elif shape_category == 'callout':
            style_parts.extend(["rounded corners", "drop shadow"])
        elif shape_category == 'rounded_rect':
            style_parts.append("rounded corners=4pt")
        elif shape_category == 'diamond':
            style_parts.append("diamond")
        elif tikz_shape_style:
            # Use the pre-computed TikZ style
            style_parts.append(tikz_shape_style)

        # Handle border/draw
        if has_border:
            if border_color:
                r, g, b = border_color
                style_parts.append(f"draw={{rgb,255:red,{r};green,{g};blue,{b}}}")
            else:
                style_parts.append("draw")

        # Handle fill color
        if fill_color:
            r, g, b, opacity = fill_color
            style_parts.append(f"fill={{rgb,255:red,{r};green,{g};blue,{b}}}")
            if opacity < 1.0:
                style_parts.append(f"fill opacity={opacity:.2f}")
        elif shape_category == 'callout':
            # Default callout fill
            style_parts.append("fill=yellow!30")
        elif shape_category in ('rectangle', 'rounded_rect', 'oval') and has_border:
            # Default fill for bordered shapes
            style_parts.append("fill=blue!10")

        # Text width based on relative width - use aspect ratio-aware textwidth
        effective_textwidth = self.get_textwidth_cm()
        text_width_cm = rel_width * effective_textwidth
        style_parts.append(f"text width={text_width_cm:.1f}cm")
        style_parts.append("align=center" if shape_category in ('oval', 'callout', 'diamond') else "align=left")
        style_parts.append("anchor=north west")

        style_str = ", ".join(style_parts)

        # Text formatting
        font_cmds = []
        if has_bold:
            font_cmds.append("\\bfseries")
        if has_italic:
            font_cmds.append("\\itshape")
        if font_size_pt:
            if font_size_pt < 8:
                font_cmds.append("\\tiny")
            elif font_size_pt < 10:
                font_cmds.append("\\scriptsize")
            elif font_size_pt < 11:
                font_cmds.append("\\footnotesize")
            elif font_size_pt < 12:
                font_cmds.append("\\small")
            elif font_size_pt > 18:
                font_cmds.append("\\Large")
            elif font_size_pt > 14:
                font_cmds.append("\\large")

        font_prefix = " ".join(font_cmds) + " " if font_cmds else ""

        # Text color
        color_prefix = ""
        if font_color_rgb:
            r, g, b = font_color_rgb
            color_prefix = f"\\color[RGB]{{{r},{g},{b}}}"

        # Position in TikZ coordinates with aspect ratio compensation
        x_pos = self.transform_x_coordinate(rel_x, canvas_width)
        y_pos = (1 - rel_y) * canvas_height  # Invert y

        lines.append(f"  \\node[{style_str}] at ({x_pos:.3f}\\textwidth, {y_pos:.3f}\\textheight) {{{color_prefix}{font_prefix}{text}}};")

    def slide_to_latex(self, slide, slide_num: int, slide_width=None, slide_height=None) -> dict:
        """Convert a single slide to LaTeX frame using comprehensive shape extraction.

        Returns dict with 'section', 'subsection', 'latex' keys for proper section management.
        """
        lines = []

        # Use the comprehensive content extraction
        content = self.get_all_shape_content(slide, slide_num, slide_width, slide_height)

        title = content['title']
        subtitle = content['subtitle']
        content_items = content['content_items']
        images = content['images']
        videos = content['videos']
        positioned_elements = content['positioned_elements']
        needs_unified_tikz = content['needs_unified_tikz']

        # If no title found, try to use first content as title
        if not title and content_items:
            title = content_items[0].get('raw_text', content_items[0]['text'])
            content_items = content_items[1:]

        # Extract sources from content items and from content dict
        content_items, sources = self.extract_sources_from_content(content_items)
        # Add any sources extracted during shape processing
        sources.extend(content.get('sources', []))

        # Estimate content density
        density = self.estimate_content_density(content_items, images, videos, positioned_elements)

        # Build the frame with secname/subsecname structure
        if density['is_text_heavy']:
            lines.append(r"\begin{frame}[shrink]{\secname\vspace{0.1cm}\\\textcolor{anthrazit!80!white}{\subsecname}}")
        else:
            lines.append(r"\begin{frame}{\secname\vspace{0.1cm}\\\textcolor{anthrazit!80!white}{\subsecname}}")

        # Add content as itemize if there are bullet points
        if content_items:
            has_bullets = any(item.get('is_bullet', False) for item in content_items)

            if has_bullets:
                strip_font_sizes = density['is_text_heavy']
                lines.append("\\begin{itemize}")
                current_level = 0
                for item in content_items:
                    level = item.get('level', 0)
                    text = item['text']

                    if strip_font_sizes:
                        text = re.sub(
                            r'\{\\(?:tiny|scriptsize|footnotesize|small|normalsize|large|Large|LARGE|huge|Huge)\s+([^{}]*(?:\{[^{}]*\}[^{}]*)*)\}',
                            r'\1',
                            text
                        )

                    while level > current_level:
                        lines.append("  " * (current_level + 1) + "\\begin{itemize}")
                        current_level += 1
                    while level < current_level:
                        lines.append("  " * current_level + "\\end{itemize}")
                        current_level -= 1

                    lines.append("  " * (level + 1) + f"\\item {text}")

                while current_level > 0:
                    lines.append("  " * current_level + "\\end{itemize}")
                    current_level -= 1
                lines.append("\\end{itemize}")
            else:
                for item in content_items:
                    text = item['text']
                    lines.append(f"{text}")
                    lines.append("")

        # Maximum image height
        max_image_height = 0.75

        # Render images and positioned elements
        if needs_unified_tikz:
            # Unified TikZ rendering for multiple images or images + shapes
            lines.extend(self.render_unified_tikz(images, positioned_elements, max_image_height))
        elif images:
            # Single image or overlapping images without positioned elements
            image_groups = self.group_overlapping_images(images)
            has_overlapping = any(len(g) > 1 for g in image_groups)

            if has_overlapping:
                for group in image_groups:
                    if len(group) == 1:
                        img_info = group[0]
                        width_ratio = img_info.get('width_ratio', 0.8)
                        lines.extend(self.render_image_with_overlays(img_info, width_ratio, max_image_height))
                    else:
                        lines.append("% Overlapping images with animation")
                        lines.append("\\begin{center}")
                        lines.append("\\begin{tikzpicture}")
                        for idx, img_info in enumerate(group):
                            filename = img_info['filename']
                            width_ratio = img_info.get('width_ratio', 0.8)
                            if idx == 0:
                                lines.append(f"  \\node[anchor=center] (img{idx}) {{\\includegraphics[width={width_ratio}\\textwidth,height={max_image_height}\\textheight,keepaspectratio]{{{filename}}}}};")
                            else:
                                overlay_num = idx + 1
                                lines.append(f"  \\onslide<{overlay_num}->{{\\node[anchor=center] at (img0.center) {{\\includegraphics[width={width_ratio}\\textwidth,height={max_image_height}\\textheight,keepaspectratio]{{{filename}}}}};}}")
                        lines.append("\\end{tikzpicture}")
                        lines.append("\\end{center}")
            else:
                for img_info in images:
                    width_ratio = img_info.get('width_ratio', 0.8)
                    if density['is_image_heavy'] and density['recommended_image_height']:
                        height_ratio = density['recommended_image_height']
                    else:
                        height_ratio = max_image_height
                    lines.extend(self.render_image_with_overlays(img_info, width_ratio, height_ratio))
        elif positioned_elements:
            # Only positioned elements, no images
            lines.append("\\begin{center}")
            lines.append("\\begin{tikzpicture}")
            for elem in positioned_elements:
                self._render_positioned_element(lines, elem, 0.9, max_image_height)
            lines.append("\\end{tikzpicture}")
            lines.append("\\end{center}")

        # Add videos
        for vid in videos:
            lines.append(f"% Video: {vid}")
            lines.append("\\includemovie[")
            lines.append("    inline=false,")
            lines.append("    attach=false,")
            lines.append("    autoplay,")
            lines.append(f"    text={{\\includegraphics[width=0.9\\textwidth,height={max_image_height}\\textheight,keepaspectratio]{{videos/generic-thumbnail.jpg}}}}")
            vid_path = "videos/" + vid
            lines.append(f"]{{0.9\\textwidth}}{{{max_image_height}\\textheight}}{{{vid_path}}}")

        # Add source command if sources were found
        if sources:
            combined_sources = "; ".join(sources)
            lines.append(f"    \\source{{{self.escape_latex(combined_sources)}}}")

        lines.append("\\end{frame}")
        lines.append("")

        return {
            'section': title if title else "Slide",
            'subsection': subtitle if subtitle else "",
            'latex': '\n'.join(lines)
        }

    def sanitize_filename(self, name: str) -> str:
        """Convert a string to a safe filename."""
        # Replace problematic characters
        safe = re.sub(r'[^\w\s\-]', '', name)
        safe = re.sub(r'\s+', '-', safe)
        return safe[:50]  # Limit length

    def convert(self) -> dict:
        """Convert the entire presentation to LaTeX.

        Returns a dict with:
        - 'title_info': dict with title, subtitle, author, date
        - 'sections': list of dicts with 'name', 'filename', 'content' (per-section LaTeX)
        - 'thank_you': dict with thank you slide info
        """
        # Patch the PPTX to handle missing content types
        patched_path = patch_pptx_content_types(str(self.input_path))
        temp_dir_to_cleanup = None
        if patched_path != str(self.input_path):
            temp_dir_to_cleanup = os.path.dirname(patched_path)

        try:
            prs = Presentation(patched_path)
        except Exception as e:
            if temp_dir_to_cleanup:
                shutil.rmtree(temp_dir_to_cleanup, ignore_errors=True)
            raise

        # Get slide dimensions for image sizing
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Detect source aspect ratio for proper coordinate transformation
        source_ratio = self.detect_aspect_ratio(slide_width, slide_height)
        ratio_name = self.get_aspect_ratio_name(source_ratio)
        print(f"  Source presentation aspect ratio: {ratio_name} ({source_ratio:.3f})")
        if abs(source_ratio - self.target_aspect_ratio) > 0.1:
            print(f"  Converting to 16:9 output - coordinates will be transformed")

        # Extract title slide info from first slide
        title_info = {'title': 'Presentation', 'subtitle': '', 'author': '', 'date': '\\today'}
        if prs.slides:
            first_slide = prs.slides[0]
            if self.is_title_slide(first_slide, 0):
                title_info.update(self.extract_title_info(first_slide))

        # Process each slide and organize by section
        sections = []  # List of {'name': section_name, 'filename': safe_filename, 'content': [lines]}
        current_section = None
        current_section_data = None
        current_subsection = None
        thank_you_info = None

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1

            # Skip title slide (already handled with \maketitle)
            if i == 0 and self.is_title_slide(slide, i):
                continue

            # Handle thank you slide specially
            if self.is_thank_you_slide(slide):
                thank_text = "Thank you for your attention!"
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        text = self.clean_text(shape.text)
                        if text and ('thank' in text.lower() or 'danke' in text.lower()):
                            thank_text = text
                            break
                thank_you_info = {'text': thank_text}
                continue

            # Get slide content with section/subsection info
            slide_data = self.slide_to_latex(slide, slide_num, slide_width, slide_height)
            section = slide_data['section']
            subsection = slide_data['subsection']
            frame_latex = slide_data['latex']

            # Check if we need to start a new section
            if section != current_section:
                # Save previous section if exists
                if current_section_data:
                    sections.append(current_section_data)

                # Start new section
                section_num = len(sections) + 1
                safe_name = self.sanitize_filename(section)
                current_section_data = {
                    'name': section,
                    'filename': f"section-{section_num:02d}-{safe_name}",
                    'content': []
                }
                current_section = section
                current_subsection = None

                # Add section command
                escaped_section = self.escape_latex(section)
                current_section_data['content'].append(f"\\section{{{escaped_section}}}")

            # Emit \subsection only when subsection changes
            if subsection != current_subsection:
                escaped_subsection = self.escape_latex(subsection) if subsection else ""
                current_section_data['content'].append(f"\\subsection{{{escaped_subsection}}}")
                current_subsection = subsection

            # Add frame content
            current_section_data['content'].append(f"% Slide {slide_num}")
            current_section_data['content'].append(frame_latex)

        # Don't forget the last section
        if current_section_data:
            sections.append(current_section_data)

        # Clean up temporary patched file
        if temp_dir_to_cleanup:
            shutil.rmtree(temp_dir_to_cleanup, ignore_errors=True)

        return {
            'title_info': title_info,
            'sections': sections,
            'thank_you': thank_you_info
        }

    def create_output_folder(self) -> Path:
        """Create the output folder structure for the converted presentation."""
        # Get presentation name without extension
        pptx_name = self.input_path.stem

        # Create output in tex-output subfolder (relative to script location)
        script_dir = Path(__file__).parent
        tex_output_dir = script_dir / 'tex-output'
        tex_output_dir.mkdir(exist_ok=True)

        # Create presentation-specific folder
        output_folder = tex_output_dir / pptx_name
        output_folder.mkdir(exist_ok=True)

        # Create subfolders
        (output_folder / 'fig').mkdir(exist_ok=True)
        (output_folder / 'videos').mkdir(exist_ok=True)
        (output_folder / 'theme').mkdir(exist_ok=True)

        # Copy generic video thumbnail to videos folder
        script_dir = Path(__file__).parent
        template_dir = script_dir / 'template'
        thumbnail_src = template_dir / 'videos' / 'generic-thumbnail.jpg'
        if thumbnail_src.exists():
            shutil.copy2(thumbnail_src, output_folder / 'videos' / 'generic-thumbnail.jpg')

        return output_folder

    def copy_theme_files(self, output_folder: Path):
        """Copy theme files to the output folder."""
        # Look for theme files in template/ subdirectory
        script_dir = Path(__file__).parent
        template_dir = script_dir / 'template'

        # Copy beamertheme.sty - try template dir first, then input parent
        theme_sty = template_dir / 'beamertheme.sty'
        if not theme_sty.exists():
            theme_sty = self.input_path.parent / 'beamertheme.sty'
        if theme_sty.exists():
            shutil.copy2(theme_sty, output_folder / 'beamertheme.sty')

        # Copy theme folder contents - try template dir first, then input parent
        theme_src = template_dir / 'theme'
        if not theme_src.exists():
            theme_src = self.input_path.parent / 'theme'
        theme_dst = output_folder / 'theme'
        if theme_src.exists():
            # Copy all files from theme folder
            for item in theme_src.iterdir():
                if item.is_file():
                    shutil.copy2(item, theme_dst / item.name)
                elif item.is_dir():
                    dst_subdir = theme_dst / item.name
                    if dst_subdir.exists():
                        shutil.rmtree(dst_subdir)
                    shutil.copytree(item, dst_subdir)

    def generate_main_tex(self, output_folder: Path, conversion_data: dict) -> str:
        """Generate the main TEX file from template."""
        # Look for template in template/ subdirectory
        script_dir = Path(__file__).parent
        template_dir = script_dir / 'template'
        template_path = template_dir / 'beamer-main-template.tex'
        if not template_path.exists():
            template_path = self.input_path.parent / 'beamer-main-template.tex'

        if template_path.exists():
            with open(template_path, 'r', encoding='utf-8') as f:
                template = f.read()
        else:
            # Fallback template
            template = """% !TeX spellcheck = en_US
\\documentclass[aspectratio=169]{beamer}
\\usepackage{beamertheme}

\\title{PRESENTATION TITLE}
\\author[AUTHOR]{AUTHOR}
\\subtitle{DATE}

\\begin{document}

\\maketitle

\\begin{frame}{Table of Contents}
    \\tableofcontents
\\end{frame}

\\include{REPLACE WITH SLIDE SECTION 1}
\\include{REPLACE WITH SLIDE SECTION 2}
\\include{REPLACE WITH SLIDE SECTION ...}

\\thankyou{Thank you for your attention}{AUTHOR}{}{EMAIL}{theme/logos/drop.png}
\\end{document}
"""

        title_info = conversion_data['title_info']
        sections = conversion_data['sections']
        thank_you = conversion_data['thank_you']

        # Replace placeholders
        main_tex = template

        # Replace \usetheme{lww} with \usepackage{beamertheme} for compatibility
        main_tex = re.sub(r'\\usetheme\{lww\}', r'\\usepackage{beamertheme}', main_tex)

        # Replace title
        main_tex = main_tex.replace('PRESENTATION TITLE', self.escape_latex(title_info['title']))

        # Replace author (handle both forms)
        author = title_info['author'] if title_info['author'] else 'Author'
        main_tex = main_tex.replace('\\author[AUTHOR]{AUHTOR}', f'\\author[{self.escape_latex(author)}]{{{self.escape_latex(author)}}}')
        main_tex = main_tex.replace('\\author[AUTHOR]{AUTHOR}', f'\\author[{self.escape_latex(author)}]{{{self.escape_latex(author)}}}')
        main_tex = re.sub(r'AUTHOR', self.escape_latex(author), main_tex)

        # Replace date/subtitle
        subtitle = title_info['subtitle'] if title_info['subtitle'] else '\\today'
        main_tex = main_tex.replace('\\subtitle{DATE}', f'\\subtitle{{{self.escape_latex(subtitle)}}}')

        # Replace email placeholder
        main_tex = main_tex.replace('EMAIL', 'email@example.com')

        # Generate include statements
        include_lines = []
        for section in sections:
            include_lines.append(f"\\include{{{section['filename']}}}")

        # Replace the template include placeholders
        # Find and replace the include block
        include_pattern = r'\\include\{REPLACE WITH SLIDE SECTION[^}]*\}(\s*\\include\{REPLACE WITH SLIDE SECTION[^}]*\})*'
        include_block = '\n'.join(include_lines)
        # Escape backslashes for regex replacement
        include_block_escaped = include_block.replace('\\', '\\\\')
        main_tex = re.sub(include_pattern, include_block_escaped, main_tex)

        # Update thank you slide
        if thank_you:
            thank_text = self.escape_latex(thank_you['text'])
        else:
            thank_text = 'Thank you for your attention!'

        # Replace the thankyou command
        thankyou_pattern = r'\\thankyou\{[^}]*\}\{[^}]*\}\{[^}]*\}\{[^}]*\}\{[^}]*\}'
        thankyou_replacement = f'\\\\thankyou{{{thank_text}}}{{{self.escape_latex(author)}}}{{}}{{email@example.com}}{{theme/logos/drop.png}}'
        main_tex = re.sub(thankyou_pattern, thankyou_replacement, main_tex)

        return main_tex

    def generate_texstudio_project(self, output_folder: Path, main_tex_name: str) -> str:
        """Generate the TexStudio project file."""
        project = {
            "InternalPDFViewer": {
                "Embedded": True,
                "File": main_tex_name.replace('.tex', '.pdf')
            },
            "Session": {
                "Bookmarks": [],
                "CurrentFile": main_tex_name,
                "FileVersion": 1,
                "Files": [
                    {
                        "Col": 0,
                        "EditorGroup": 0,
                        "FileName": main_tex_name,
                        "FirstLine": 0,
                        "FoldedLines": "",
                        "Line": 0
                    }
                ],
                "MasterFile": "",
                "VerticalSplit": False
            }
        }
        return json.dumps(project, indent=4)

    def save(self):
        """Convert and save the LaTeX output to organized folder structure."""
        # Create output folder structure FIRST (before conversion)
        output_folder = self.create_output_folder()

        # Update fig_dir and video_dir to point to new location BEFORE conversion
        self.fig_dir = output_folder / 'fig'
        self.video_dir = output_folder / 'videos'

        # Ensure directories exist
        self.fig_dir.mkdir(parents=True, exist_ok=True)
        self.video_dir.mkdir(parents=True, exist_ok=True)

        # Reset counters for fresh extraction
        self.image_counter = 0
        self.video_counter = 0

        # Now convert the presentation (images will be extracted to new location)
        conversion_data = self.convert()

        # Copy theme files
        self.copy_theme_files(output_folder)

        # Write per-section TEX files
        for section in conversion_data['sections']:
            section_path = output_folder / f"{section['filename']}.tex"
            with open(section_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(section['content']))

        # Generate and write main TEX file
        main_tex_content = self.generate_main_tex(output_folder, conversion_data)
        main_tex_path = output_folder / 'beamer-main.tex'
        with open(main_tex_path, 'w', encoding='utf-8') as f:
            f.write(main_tex_content)

        # Generate and write TexStudio project file
        project_content = self.generate_texstudio_project(output_folder, 'beamer-main.tex')
        project_path = output_folder / 'beamerProject-TexStudio.txss2'
        with open(project_path, 'w', encoding='utf-8') as f:
            f.write(project_content)

        print(f"  Output folder: {output_folder}")
        print(f"  Main file: beamer-main.tex")
        print(f"  Section files: {len(conversion_data['sections'])}")

        print(f"Converted: {self.input_path} -> {self.output_path}")
        print(f"  Images saved to: {self.fig_dir}/")
        print(f"  Videos saved to: {self.video_dir}/")


def main():
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint presentations to LaTeX Beamer format",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s                           # Convert all .pptx in pptx-input/ to beamer-slides.tex
  %(prog)s presentation.pptx         # Convert specific file to beamer-slides.tex
  %(prog)s input.pptx output.tex     # Specify custom output filename
        """
    )
    parser.add_argument(
        'input',
        nargs='?',
        help='Input PPTX file (default: process all in pptx-input/)'
    )
    parser.add_argument(
        'output',
        nargs='?',
        help='Output TEX file (default: beamer-slides.tex)'
    )
    parser.add_argument(
        '--fig-dir',
        default='fig',
        help='Directory for extracted images (default: fig)'
    )
    parser.add_argument(
        '--video-dir',
        default='videos',
        help='Directory for extracted videos (default: videos)'
    )

    args = parser.parse_args()

    # Determine input files
    if args.input:
        input_files = [Path(args.input)]
    else:
        # Process all PPTX files in pptx-input/
        input_dir = Path('pptx-input')
        if not input_dir.exists():
            print(f"Error: Input directory '{input_dir}' does not exist")
            sys.exit(1)

        input_files = list(input_dir.glob('*.pptx'))
        if not input_files:
            print(f"No .pptx files found in '{input_dir}'")
            sys.exit(1)

    # Process each file
    for input_path in input_files:
        if not input_path.exists():
            print(f"Error: File '{input_path}' does not exist")
            continue

        # Determine output path - always use beamer-slides.tex for TexStudio compatibility
        if args.output and len(input_files) == 1:
            output_path = Path(args.output)
        else:
            output_path = Path('beamer-slides.tex')

        try:
            converter = PPTXToLatexConverter(
                str(input_path),
                str(output_path),
                fig_dir=args.fig_dir,
                video_dir=args.video_dir
            )
            converter.save()
        except Exception as e:
            print(f"Error converting {input_path}: {e}")
            raise


if __name__ == '__main__':
    main()